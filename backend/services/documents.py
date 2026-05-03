"""
backend/services/documents.py
Multi-format text extraction + structured marketing analysis.
PDF / DOCX / TXT / MD / PPTX / images (Gemini Vision).
"""
from __future__ import annotations
import io
import json
import re
import base64
from typing import Optional


def _safe_decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try: return data.decode(enc)
        except Exception: continue
    return data.decode("utf-8", errors="ignore")


async def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    name = (filename or "").lower()
    if name.endswith(".pdf"):
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            n = len(reader.pages)
            if n <= 20:
                pages = reader.pages
            else:
                pages = list(reader.pages[:5]) + list(reader.pages[-5:])
            chunks = []
            for p in pages:
                try: chunks.append(p.extract_text() or "")
                except Exception: continue
            return "\n\n".join(chunks).strip()
        except Exception as e:
            return f"[pdf extraction failed: {e}]"
    if name.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            return f"[docx extraction failed: {e}]"
    if name.endswith(".pptx"):
        try:
            from pptx import Presentation  # python-pptx
            prs = Presentation(io.BytesIO(file_bytes))
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            return "\n".join(chunks)
        except Exception as e:
            return f"[pptx extraction failed: {e}]"
    if name.endswith((".txt", ".md", ".csv", ".json")):
        return _safe_decode(file_bytes)[:200000]
    if name.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif")):
        # OCR via Gemini Vision
        try:
            from ..brain.ai_router import get_router
            router = get_router()
            return await router.gemini_vision_call(
                "Transcribe ALL visible text. Then in 2 lines describe the visual.",
                file_bytes,
            )
        except Exception as e:
            return f"[image OCR failed: {e}]"
    # Fallback: best-effort decode
    return _safe_decode(file_bytes)[:50000]


async def _detect_type(text: str) -> str:
    from ..brain.ai_router import get_router
    router = get_router()
    prompt = (
        "Classify this document in ONE WORD from: pitch_deck, business_plan, "
        "marketing_report, product_spec, financial, legal, resume, article, other.\n\n"
        + text[:2000]
    )
    try:
        return (await router.simple_groq(prompt))[:40].strip().lower()
    except Exception:
        return "other"


ANALYSIS_TEMPLATE = """You are Friday, {user_name}'s personal CMO. Analyze this document
and produce a STRUCTURED marketing brief.

DOCUMENT TYPE (auto-detected): {doc_type}
USER BUSINESS: {biz}
FRIEND CONTEXT: {friend}

DOCUMENT CONTENT:
{content}

Return your analysis in EXACTLY this format (use the headers verbatim):

WHAT THIS IS: <1-2 sentences>

CORE OPPORTUNITY: <the single biggest marketing opportunity in 1-2 sentences>

TARGET AUDIENCE: <who this is for, specific>

STRENGTHS:
1. <strength>
2. <strength>
3. <strength>

GAPS:
1. <gap>
2. <gap>
3. <gap>

COMPETITIVE POSITIONING: <1 paragraph>

CHANNEL STRATEGY:
1. <channel + why>
2. <channel + why>
3. <channel + why>

ACTION PLAN:
- Phase 1 (week 1-2): <2-3 actions>
- Phase 2 (week 3-6): <2-3 actions>
- Phase 3 (month 2-3): <2-3 actions>

ONE THING: <the ONE highest-leverage move he should do this week>

SPOKEN BRIEF: <60-word friend-style spoken summary, no headers, ready for TTS>
"""


async def analyze_document_for_marketing(file_bytes: bytes, filename: str,
                                         user_context: Optional[dict] = None,
                                         friend_context: str = "",
                                         user_name: str = "Mr Vijesh"
                                         ) -> dict:
    text = await extract_text_from_file(file_bytes, filename)
    if len(text) < 30:
        return {"error": "Could not extract usable text", "filename": filename}
    doc_type = await _detect_type(text)
    from ..brain.ai_router import get_router
    router = get_router()
    biz = json.dumps(user_context or {}, default=str)[:1200]
    prompt = ANALYSIS_TEMPLATE.format(
        user_name=user_name, doc_type=doc_type, biz=biz,
        friend=friend_context[:1200], content=text[:8000],
    )
    analysis = ""
    for _call in (
        lambda: router.simple_gemini(prompt, model="gemini-2.0-flash"),
        lambda: router.simple_deepseek(prompt),
        lambda: router.simple_groq(prompt),
    ):
        try:
            analysis = await _call()
            if analysis: break
        except Exception:
            continue
    spoken = ""
    m = re.search(r"SPOKEN BRIEF:\s*(.+)", analysis, re.IGNORECASE | re.DOTALL)
    if m: spoken = m.group(1).strip()[:600]
    return {
        "filename": filename,
        "doc_type": doc_type,
        "analysis": analysis,
        "spoken": spoken,
        "extract_preview": text[:500],
    }


# Backwards-compat shim used by older code paths
async def extract_text(file_bytes: bytes, filename: str) -> str:
    return await extract_text_from_file(file_bytes, filename)


async def analyze_for_marketing(file_bytes: bytes, filename: str,
                                user_context: Optional[dict] = None) -> dict:
    return await analyze_document_for_marketing(file_bytes, filename,
                                                user_context or {}, "")
